#!/usr/bin/env python3
import sys
import os
import subprocess

# Ensure python-pptx is installed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx is not installed. Installing it now...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    # Initialize Presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # -------------------------------------------------------------
    # Premium Palette Configuration (Modern Developer Dark Mode)
    # -------------------------------------------------------------
    BG_COLOR = RGBColor(18, 18, 20)          # #121214 Deep Charcoal
    PANEL_BG = RGBColor(30, 30, 36)          # #1E1E24 Slightly lighter Gray
    TEXT_TITLE = RGBColor(248, 250, 252)     # #F8FAFC Off-white
    TEXT_BODY = RGBColor(161, 161, 170)      # #A1A1AA Soft Gray
    ACCENT_INDIGO = RGBColor(99, 102, 241)   # #6366F1 Royal Indigo
    ACCENT_TEAL = RGBColor(20, 184, 166)     # #14B8A6 Teal
    ACCENT_ORANGE = RGBColor(245, 158, 11)   # #F59E0B Warm Amber
    ACCENT_PURPLE = RGBColor(168, 85, 247)   # #A855F7 Purple

    # Helper function to apply dark background to slide
    def apply_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    # Helper function to write standardized titles
    def add_slide_header(slide, title_text, category="MACHINE LEARNING FUNDAMENTALS"):
        # Category breadcrumb
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.3))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_top = cat_tf.margin_bottom = cat_tf.margin_left = cat_tf.margin_right = 0
        p_cat = cat_tf.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = "Arial"
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_TEAL
        
        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_top = title_tf.margin_bottom = title_tf.margin_left = title_tf.margin_right = 0
        p_title = title_tf.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Arial"
        p_title.font.size = Pt(30)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_TITLE

    # Helper to add standard panels
    def draw_panel(slide, left, top, width, height, bg_color=PANEL_BG, border_color=ACCENT_INDIGO):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    # Helper function to scale and center an image on the right
    def add_right_image(slide, img_path, left=Inches(7.2), top=Inches(1.8), max_width=Inches(5.3), max_height=Inches(4.8)):
        if not os.path.exists(img_path):
            print(f"Warning: Image not found at '{img_path}'. Creating placeholder.")
            # Draw placeholder instead
            draw_panel(slide, left, top, max_width, max_height, bg_color=PANEL_BG, border_color=ACCENT_ORANGE)
            box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.8), max_width - Inches(0.4), Inches(1.5))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"Image File Missing:\n{os.path.basename(img_path)}"
            p.font.name = "Arial"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = ACCENT_ORANGE
            p.alignment = PP_ALIGN.CENTER
            return None

        # Add image temporarily to inspect size
        pic = slide.shapes.add_picture(img_path, left, top)
        aspect_ratio = pic.width / pic.height
        
        # Calculate optimal size preserving aspect ratio
        if max_width / max_height > aspect_ratio:
            # Height is the limiting factor
            pic.height = int(max_height)
            pic.width = int(max_height * aspect_ratio)
            pic.left = int(left + (max_width - pic.width) / 2)
            pic.top = int(top)
        else:
            # Width is the limiting factor
            pic.width = int(max_width)
            pic.height = int(max_width / aspect_ratio)
            pic.left = int(left)
            pic.top = int(top + (max_height - pic.height) / 2)
        return pic

    # Helper function to generate standard layout slides (Left text, Right image)
    def create_standard_slide(title, category, panel_title, bullets, img_path, border_color=ACCENT_INDIGO):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_bg(slide)
        add_slide_header(slide, title, category)
        
        # Left Panel Background
        draw_panel(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.8), border_color=border_color)
        
        # Text Frame
        box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.6), Inches(4.4))
        tf = box.text_frame
        tf.word_wrap = True
        
        # Panel Title
        p_title = tf.paragraphs[0]
        p_title.text = panel_title
        p_title.font.name = "Arial"
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = ACCENT_TEAL
        p_title.space_after = Pt(12)
        
        # Bullets formatting
        for b in bullets:
            p = tf.add_paragraph()
            p.font.name = "Arial"
            p.space_after = Pt(8)
            
            if b.startswith("  -") or b.startswith("    "):
                p.text = "    " + b.strip()
                p.font.size = Pt(12)
                p.font.color.rgb = TEXT_BODY
            elif b.startswith("-") or b.startswith("•"):
                p.text = "• " + b.lstrip("-• ").strip()
                p.font.size = Pt(13)
                p.font.color.rgb = TEXT_BODY
            else:
                p.text = b
                p.font.size = Pt(13)
                p.font.color.rgb = TEXT_BODY
                
        # Right Image Placement
        add_right_image(slide, img_path, left=Inches(7.2), top=Inches(1.8), max_width=Inches(5.3), max_height=Inches(4.8))
        return slide

    # =============================================================
    # Slide 1: Title Slide (Sleek Landing)
    # =============================================================
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    
    highlight = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = ACCENT_INDIGO
    highlight.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.5), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Fundamentals of\nMachine Learning"
    p1.font.name = "Arial"
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_TITLE
    p1.space_after = Pt(20)
    
    p2 = tf.add_paragraph()
    p2.text = "A Step-by-Step Teaching Curriculum & Core Models Guide"
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT_TEAL
    p2.space_after = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = "Instructor Slide Deck  •  Supervised, Unsupervised, Ensembles & LLMs"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_BODY

    # =============================================================
    # Slide 2: The Machine Learning Paradigm Shift
    # =============================================================
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_slide_header(slide, "The Paradigm Shift: Rules vs. Data", "MODULE 1: INTRODUCTION TO ML")
    
    # Traditional Panel
    draw_panel(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
    box_trad = slide.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(4.9), Inches(4.4))
    tf_trad = box_trad.text_frame
    tf_trad.word_wrap = True
    
    p = tf_trad.paragraphs[0]
    p.text = "Traditional Software Engineering"
    p.font.name = "Arial"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(14)
    
    bullets_trad = [
        "Human developer codes logical rules explicitly.",
        "System applies rules directly to inputs to yield answers.",
        "Highly deterministic but difficult to scale for complex patterns like vision or natural language.",
        "Example:\n  If 'free' and 'money' in text: mark as spam."
    ]
    for b in bullets_trad:
        p = tf_trad.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_BODY
        p.space_after = Pt(10)
        
    # Machine Learning Panel
    draw_panel(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.8))
    box_ml = slide.shapes.add_textbox(Inches(7.3), Inches(2.0), Inches(4.9), Inches(4.4))
    tf_ml = box_ml.text_frame
    tf_ml.word_wrap = True
    
    p = tf_ml.paragraphs[0]
    p.text = "Machine Learning Paradigm"
    p.font.name = "Arial"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_INDIGO
    p.space_after = Pt(14)
    
    bullets_ml = [
        "Algorithm accepts raw data paired with historical outcomes.",
        "Model self-optimizes internal parameters to map X to y.",
        "Excels at multi-variable high-dimensional datasets that defy simple rule creation.",
        "Example:\n  Analyze millions of emails to extract probabilistic spam indicators automatically."
    ]
    for b in bullets_ml:
        p = tf_ml.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_BODY
        p.space_after = Pt(10)

    # =============================================================
    # Slide 3: The 4 Learning Paradigms
    # =============================================================
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_slide_header(slide, "The Four Machine Learning Paradigms", "MODULE 2: LEARNING PARADIGMS")
    
    paradigms = [
        ("Supervised", "Labeled data drives training. Target labels (y) are paired with features (X).", "Credit Scoring, Pricing", ACCENT_INDIGO, Inches(0.8)),
        ("Unsupervised", "Unlabeled data. Discovers hidden structural groupings directly from features (X).", "Customer Clustering", ACCENT_TEAL, Inches(3.85)),
        ("Semi-Supervised", "Leverages small labeled dataset + huge unlabeled pool to reduce annotation costs.", "Image Annotation", ACCENT_ORANGE, Inches(6.9)),
        ("Reinforcement", "Interactive trial-and-error. Agent optimizes policy using environment rewards.", "Autonomous Driving", ACCENT_PURPLE, Inches(9.95))
    ]
    
    for title, desc, eg, color, left in paradigms:
        draw_panel(slide, left, Inches(1.8), Inches(2.6), Inches(4.8), border_color=color)
        
        box = slide.shapes.add_textbox(left + Inches(0.15), Inches(2.0), Inches(2.3), Inches(4.4))
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(12)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_BODY
        p.space_after = Pt(24)
        
        p = tf.add_paragraph()
        p.text = "Example Use Case:"
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_TITLE
        
        p = tf.add_paragraph()
        p.text = eg
        p.font.size = Pt(13)
        p.font.italic = True
        p.font.color.rgb = color

    # =============================================================
    # Slide 4: Supervised Regression
    # =============================================================
    create_standard_slide(
        title="Regression: Predicting Continuous Targets",
        category="MODULE 3: SUPERVISED REGRESSION",
        panel_title="Continuous Model Formulation",
        bullets=[
            "Ordinary Least Squares (OLS): Fits a linear equation: y_hat = theta_0 + theta_1 * x_1 + ...",
            "Cost Function: Minimizes Mean Squared Error (MSE) to align the regression line.",
            "L2 Regularization (Ridge): Adds penalty lambda * sum(theta^2). Shrinks weights close to zero, decreasing variance.",
            "L1 Regularization (Lasso): Adds penalty lambda * sum(|theta|). Drives coefficients to exactly zero (Soft-Thresholding), doing automatic feature selection.",
            "Visual: Framework OLS vs. Ridge vs. Lasso on housing per-sq-ft."
        ],
        img_path="plots/examples_1_regression.png",
        border_color=ACCENT_TEAL
    )

    # =============================================================
    # Slide 5: Supervised Classification
    # =============================================================
    create_standard_slide(
        title="Classification: Categorizing Observations",
        category="MODULE 4: SUPERVISED CLASSIFICATION",
        panel_title="Decision Boundary Approaches",
        bullets=[
            "Logistic Regression: Maps inputs to probabilities using the Sigmoid function: 1 / (1 + e^-z).",
            "K-Nearest Neighbors (KNN): Non-parametric model. Classifies a query point based on distance-weighted votes of its K closest neighbors.",
            "Support Vector Machines (SVM): Finds a hyperplane that maximizes the margin between classes. Employs the Kernel Trick for non-linear boundaries.",
            "Naive Bayes: Probabilistic classifier based on Bayes Theorem. Assumes strong conditional independence among features.",
            "Visual: Multi-class decision boundaries of Scikit-Learn models on the Iris dataset."
        ],
        img_path="plots/classification_decision_boundaries.png",
        border_color=ACCENT_INDIGO
    )

    # =============================================================
    # Slide 6: Distance Metrics in Machine Learning
    # =============================================================
    create_standard_slide(
        title="Distance Metrics & Vector Spatial Relations",
        category="MODULE 4: SUPERVISED CLASSIFICATION",
        panel_title="Quantifying Geometric Closeness",
        bullets=[
            "Euclidean Distance: Straight-line (L2 norm) distance. Sensitive to large coordinate deviations.",
            "Manhattan Distance: Grid-based (L1 norm) distance. Sum of absolute differences along coordinate axes.",
            "Cosine Similarity: Measures the angular alignment between two vectors. Independent of vector magnitude (widely used in text/embeddings).",
            "Chebyshev Distance: Maximum absolute coordinate difference (L-infinity norm).",
            "Visual: Distance metric bounds (Circle, Diamond, Square) in a 2D coordinate space."
        ],
        img_path="plots/distance_metrics.png",
        border_color=ACCENT_ORANGE
    )

    # =============================================================
    # Slide 7: Decision Trees & Random Forests
    # =============================================================
    create_standard_slide(
        title="Tree-Based Models & Bagging Ensembles",
        category="MODULE 4: SUPERVISED CLASSIFICATION",
        panel_title="Recursive Splits & Forest Aggregation",
        bullets=[
            "Decision Trees: Splits features recursively to maximize sample purity.",
            "Purity Metrics: Evaluated using Gini Impurity or Entropy (Information Gain). Highly prone to overfitting.",
            "Random Forest Ensemble: Trains multiple independent trees in parallel.",
            "Bootstrapping: Each tree is trained on a random subset of rows (drawn with replacement).",
            "Feature Subspacing: Selects a random subset of features (e.g. sqrt(D)) at each split to decorrelate tree errors.",
            "Visual: Decision Tree splits vs. Random Forest averaged boundary."
        ],
        img_path="plots/examples_4_decision_tree.png",
        border_color=ACCENT_TEAL
    )

    # =============================================================
    # Slide 8: Ensemble Boosting & XGBoost
    # =============================================================
    create_standard_slide(
        title="Boosting Paradigms & XGBoost",
        category="MODULE 5: ENSEMBLE LEARNING & XGBOOST",
        panel_title="Sequential Optimization",
        bullets=[
            "Sequential Learning: Weak learners (e.g. shallow decision trees) are trained sequentially rather than in parallel.",
            "Residual Fitting: Each successive tree is trained to predict the residual errors (gradients) of the cumulative ensemble.",
            "XGBoost (Extreme Gradient Boosting): Highly optimized library.",
            "Regularized Objective: Incorporates L1 & L2 parameter constraints directly into the split objective.",
            "Implementation: Employs parallel split searches, block structures, and handles missing/sparse data automatically."
        ],
        img_path="plots/bagging_vs_boosting.png",
        border_color=ACCENT_PURPLE
    )

    # =============================================================
    # Slide 9: Unsupervised Learning - Clustering
    # =============================================================
    create_standard_slide(
        title="Clustering: Grouping Unlabeled Points",
        category="MODULE 6: UNSUPERVISED LEARNING",
        panel_title="Unsupervised Spatial Grouping",
        bullets=[
            "K-Means: Partitions data into K clusters. Iteratively updates cluster centroids to minimize within-cluster sum of squares (Inertia).",
            "Hierarchical Clustering: Builds nested tree structures (dendrograms) via agglomerative (bottom-up) merges.",
            "DBSCAN: Density-based algorithm. Groups core points within radius Eps having min_samples. Isolates outliers as noise.",
            "Visual: Scikit-Learn K-Means customer segmentation showing centroids and color-mapped clusters."
        ],
        img_path="plots/examples_6_kmeans.png",
        border_color=ACCENT_TEAL
    )

    # =============================================================
    # Slide 10: Model Evaluation & Validation
    # =============================================================
    create_standard_slide(
        title="Model Evaluation & Generalization",
        category="MODULE 7: EVALUATION & VALIDATION",
        panel_title="Validation & Performance Metrics",
        bullets=[
            "Bias-Variance Tradeoff: Underfitting (High Bias) vs. Overfitting (High Variance). Goal is minimizing total test error.",
            "K-Fold Cross-Validation: Splits data into K folds, cycling train/test roles to obtain robust performance bounds.",
            "Classification Metrics: Derived from the Confusion Matrix.",
            "  - Precision: TP / (TP + FP) (Quality of positive predictions).",
            "  - Recall: TP / (TP + FN) (Coverage of actual positive samples).",
            "  - F1-Score: Harmonic mean of Precision and Recall.",
            "  - ROC-AUC: True Positive vs. False Positive rate probability area."
        ],
        img_path="plots/examples_8_evaluation.png",
        border_color=ACCENT_INDIGO
    )

    # =============================================================
    # Slide 11: Module 8: Code Walkthroughs
    # =============================================================
    create_standard_slide(
        title="Framework vs. From-Scratch Algorithms",
        category="MODULE 8: CODE WALKTHROUGHS",
        panel_title="Pedagogical Implementation Strategy",
        bullets=[
            "Dual-Path Curriculum: Every algorithm is explored via two implementations: Scikit-Learn and Pure Python.",
            "Framework Path: Teaches APIs, hyperparameter tuning, pipelines, and industry best practices.",
            "From-Scratch Path: Avoids libraries. Uses standard lists and math functions to write loops, gradients, and Gini calculations.",
            "De-mystifying the 'Black Box': Solidifies mathematical equations directly into concrete, readable code.",
            "Visual: Pure Python Gradient Descent line-fitting convergence over epochs."
        ],
        img_path="plots/scratch_1_regression.png",
        border_color=ACCENT_ORANGE
    )

    # =============================================================
    # Slide 12: The Artificial Neuron (Perceptron)
    # =============================================================
    create_standard_slide(
        title="The Artificial Neuron (Perceptron)",
        category="MODULE 9: DEEP LEARNING FOUNDATIONS",
        panel_title="Single Neuron Computation",
        bullets=[
            "Perceptron Concept: Mimics a biological neuron by summing weighted inputs and passing the result to an activation function.",
            "Mathematical Formulation:",
            "  - Pre-activation sum: z = Σ (w_i * x_i) + b = w^T * x + b",
            "  - Activation output: y = φ(z) = φ(w^T * x + b)",
            "  - Where w represents connection weights, and b is the bias.",
            "AND Logic Gate Example (Linear Decision Boundary):",
            "  - Parameters: w_1 = 1.0, w_2 = 1.0, b = -1.5",
            "  - Activation φ(z): Step function (1 if z >= 0, else 0)",
            "  - Query [0, 0] -> z = -1.5 -> y = 0",
            "  - Query [1, 0] -> z = -0.5 -> y = 0",
            "  - Query [1, 1] -> z = +0.5 -> y = 1 (AND logic holds)"
        ],
        img_path="plots/perceptron_diagram.png",
        border_color=ACCENT_INDIGO
    )

    # =============================================================
    # Slide 13: Activation Functions & Non-Linearity
    # =============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide)
    add_slide_header(slide, "Activation Functions & Non-Linearity", "MODULE 9: DEEP LEARNING FOUNDATIONS")
    
    # 5 Activation Column Panels
    activations = [
        ("Sigmoid", "φ(z) = 1 / (1 + e^-z)", "Range: (0, 1)\n\nφ'(z) = φ(z)(1 - φ(z))\n\nOutputs represent probabilities.\n\nCons: Vanishing gradients.", ACCENT_INDIGO, Inches(0.8)),
        ("Tanh", "φ(z) = tanh(z)", "Range: (-1, 1)\n\nφ'(z) = 1 - φ(z)^2\n\nZero-centered outputs help stabilize deep networks.\n\nCons: Vanishing gradients.", ACCENT_TEAL, Inches(3.2)),
        ("ReLU", "φ(z) = max(0, z)", "Range: [0, ∞)\n\nφ'(z) = 1 (z > 0), else 0\n\nComputationally cheap; resolves vanishing gradient.\n\nCons: Dying ReLU problem.", ACCENT_ORANGE, Inches(5.6)),
        ("Leaky ReLU", "φ(z) = max(αz, z)", "Range: (-∞, ∞)\n\nφ'(z) = 1 (z > 0), else α\n\nPrevents dead neurons by keeping a small gradient.\n\nCons: Extra parameter α.", ACCENT_PURPLE, Inches(8.0)),
        ("Softmax", "P_i = e^z_i / Σ e^z_j", "Range: (0, 1) (sums to 1)\n\nDerivative: P_i(δ_ij - P_j)\n\nOutputs normalized probabilities for classification.\n\nUsage: Final output layer.", ACCENT_TEAL, Inches(10.4))
    ]
    
    for title, formula, details, color, left in activations:
        draw_panel(slide, left, Inches(1.8), Inches(2.13), Inches(4.8), border_color=color)
        
        box = slide.shapes.add_textbox(left + Inches(0.1), Inches(2.0), Inches(1.93), Inches(4.4))
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(10)
        
        p = tf.add_paragraph()
        p.text = formula
        p.font.name = "Courier New"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TEXT_TITLE
        p.space_after = Pt(14)
        
        for line in details.split("\n\n"):
            p = tf.add_paragraph()
            p.text = "• " + line
            p.font.name = "Arial"
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_BODY
            p.space_after = Pt(8)

    # =============================================================
    # Slide 14: Deep Neural Networks (ANN / DNN)
    # =============================================================
    create_standard_slide(
        title="Deep Neural Networks (ANN / DNN)",
        category="MODULE 9: DEEP LEARNING FOUNDATIONS",
        panel_title="Layered Feedforward Architecture",
        bullets=[
            "Multi-Layer Perceptron (MLP): Stacks layers of neurons together. Information flows forward from input, through hidden layers, to output.",
            "Forward Propagation Equations (Layer l):",
            "  - Pre-activation: z^[l] = W^[l] * a^[l-1] + b^[l]",
            "  - Activation: a^[l] = g^[l](z^[l])",
            "  - Where a^[0] = x (the raw input features).",
            "  - W^[l] is the weight matrix of shape (n^[l] x n^[l-1]).",
            "  - b^[l] is the bias vector of shape (n^[l] x 1).",
            "  - g^[l] is the layer's activation function.",
            "Universal Approximation Theorem: A feedforward network with a single hidden layer and non-linear activations can approximate any continuous function."
        ],
        img_path="plots/dnn_architecture.png",
        border_color=ACCENT_TEAL
    )

    # =============================================================
    # Slide 15: Backpropagation (Chain Rule)
    # =============================================================
    create_standard_slide(
        title="Backpropagation: Gradient Flow via Chain Rule",
        category="MODULE 9: DEEP LEARNING FOUNDATIONS",
        panel_title="Backward Error Propagation",
        bullets=[
            "Training Objective: Find weights and biases that minimize a loss function L(y, y_hat).",
            "The Chain Rule: Calculates the sensitivity of loss to a specific weight:",
            "  - ∂L / ∂w_ij^[l] = (∂L / ∂a_i^[l]) * (∂a_i^[l] / ∂z_i^[l]) * (∂z_i^[l] / ∂w_ij^[l])",
            "Error Term definition: δ_i^[l] = ∂L / ∂z_i^[l]",
            "1. Output Layer Error (Layer L):",
            "  - δ_i^[L] = (∂L / ∂a_i^[L]) * g'^[L](z_i^[L])",
            "2. Hidden Layer Error (Layer l) propagated backward:",
            "  - δ_j^[l] = (Σ (δ_k^[l+1] * w_kj^[l+1])) * g'^[l](z_j^[l])",
            "3. Gradients calculation:",
            "  - ∂L / ∂w_ji^[l] = δ_j^[l] * a_i^[l-1]",
            "  - ∂L / ∂b_j^[l] = δ_j^[l]"
        ],
        img_path="plots/backpropagation_diagram.png",
        border_color=ACCENT_PURPLE
    )

    # =============================================================
    # Slide 16: Convolutional Neural Networks (CNN)
    # =============================================================
    create_standard_slide(
        title="Convolutional Neural Networks (CNN)",
        category="MODULE 9: DEEP LEARNING FOUNDATIONS",
        panel_title="Spatial Features & Translation Invariance",
        bullets=[
            "Overfitting Challenge: Fully connected networks do not scale well to images. A 1000x1000 RGB image requires 3M inputs per neuron.",
            "CNN Architectural Solutions:",
            "  1. Local Connectivity: Neurons connect only to small local patches.",
            "  2. Shared Weights: Kernels scan the entire input, sharing parameters.",
            "Output Size Formula for Convolutional Layer:",
            "  - O = floor((W - K + 2P) / S) + 1",
            "  - W = input size, K = kernel size, P = padding, S = stride.",
            "Pooling Layer: Reduces representation size and introduces invariance.",
            "  - Max Pooling: Extracts the maximum value in a window.",
            "  - Average Pooling: Computes the mean of the window."
        ],
        img_path="plots/cnn_architecture.png",
        border_color=ACCENT_ORANGE
    )

    # =============================================================
    # Slide 17: Module 10.1: Tokenization (BPE)
    # =============================================================
    create_standard_slide(
        title="Tokenization: Byte Pair Encoding (BPE)",
        category="MODULE 10: LLMS & TRANSFORMERS",
        panel_title="Subword Text Segmentation",
        bullets=[
            "Flaws of Word/Char Tokenizers: Massive vocabulary sizes (memory heavy) or inability to handle unseen out-of-vocabulary (OOV) words.",
            "BPE Concept: Dynamically merges the most frequent adjacent byte/character pairs into new subword units.",
            "Algorithm Steps:",
            "  1. Initialize vocabulary with individual characters + '</w>' marker.",
            "  2. Count frequencies of adjacent token pairs (bigrams) in the corpus.",
            "  3. Merge the most frequent bigram and add it to the vocabulary.",
            "  4. Repeat for N merge iterations.",
            "Visual: Vocabulary size growth vs. BPE merge rules application."
        ],
        img_path="plots/llm_1_bpe_vocab.png",
        border_color=ACCENT_TEAL
    )

    # =============================================================
    # Slide 18: Module 10.2: Vector Semantics (Word2Vec)
    # =============================================================
    create_standard_slide(
        title="Vector Semantics: Word Embeddings",
        category="MODULE 10: LLMS & TRANSFORMERS",
        panel_title="Continuous Word Semantics",
        bullets=[
            "Embedding Concept: Maps raw strings to dense, low-dimensional continuous vectors where geometric closeness represents semantic meaning.",
            "Word2Vec Skip-gram: Neural network trained to predict context words given a target input word.",
            "Negative Sampling Loss: Avoids expensive full softmax by training a binary classifier on positive context pairs vs. random negative pairs.",
            "Optimization: Trained using the Adam optimizer implemented from scratch. Uses adaptive learning rates for faster, more stable coordinate updates.",
            "Visual: Trained 2D Word2Vec scatter plot showing semantic clusters."
        ],
        img_path="plots/llm_2_word2vec.png",
        border_color=ACCENT_INDIGO
    )

    # =============================================================
    # Slide 19: Module 10.3: Transformer Network Architecture
    # =============================================================
    create_standard_slide(
        title="The Transformer Network Architecture",
        category="MODULE 10: LLMS & TRANSFORMERS",
        panel_title="Attention-Based Sequence Processing",
        bullets=[
            "Parallelization Breakthrough: Eliminates recurrence loops of RNNs/LSTMs. Processes all tokens simultaneously.",
            "Core Encoder-Decoder Structure:",
            "  - Encoder: Learns bidirectional contextual representations of the input sequence.",
            "  - Decoder: Autoregressively generates output tokens, querying encoder keys and values.",
            "Key Mechanisms:",
            "  - Multi-Head Attention (MHA): Integrates multiple parallel self-attention views to capture complex context relations.",
            "  - Residual Connections & LayerNorm: Stabilize gradient flow in deep stacks.",
            "  - Feed-Forward Networks (FFN): Applied pointwise to each token position."
        ],
        img_path="plots/transformer_architecture.png",
        border_color=ACCENT_PURPLE
    )

    # =============================================================
    # Slide 20: Module 10.4: Transformer Encoder
    # =============================================================
    create_standard_slide(
        title="The Encoder: Self-Attention & Positional Encoding",
        category="MODULE 10: LLMS & TRANSFORMERS",
        panel_title="Parallel Contextual Encoding",
        bullets=[
            "Positional Encoding: Adds sine/cosine wave patterns to input embeddings to preserve token sequence order.",
            "Scaled Dot-Product Self-Attention: Projects input matrices into Query (Q), Key (K), and Value (V) representations.",
            "Alignment Formula: Attention(Q, K, V) = Softmax(Q K^T / sqrt(d_k)) V. Scales by sqrt(d_k) to prevent gradient vanishing.",
            "Context Resolution: Dynamically shifts a word's vector representation based on surrounding tokens (e.g. 'bank' of river vs. money).",
            "Visual: Sinusoidal Positional Encoding matrix heatmap."
        ],
        img_path="plots/llm_3_positional_encoding.png",
        border_color=ACCENT_PURPLE
    )

    # =============================================================
    # Slide 21: Module 10.5: Transformer Decoder
    # =============================================================
    create_standard_slide(
        title="The Decoder: Causal Masking & Cross-Attention",
        category="MODULE 10: LLMS & TRANSFORMERS",
        panel_title="Autoregressive Target Generation",
        bullets=[
            "Autoregressive Generation: Predicts tokens sequentially, appending prior outputs as inputs for successive steps.",
            "Causal Masking: Adds a lower-triangular mask (-infinity for future positions) to attention scores prior to softmax.",
            "Look-Ahead Prevention: Mathematically sets attention weights to future tokens to exactly 0, preventing future-leakage.",
            "Encoder-Decoder Cross-Attention: Decoder queries (Q) attend to Encoder keys (K) and values (V), linking target to source.",
            "Visual: Decoder Causal Mask lower-triangular attention boundaries."
        ],
        img_path="plots/llm_4_causal_mask.png",
        border_color=ACCENT_ORANGE
    )

    # =============================================================
    # Slide 22: Module 10.6: End-to-End LLMSeq2Seq
    # =============================================================
    create_standard_slide(
        title="End-to-End Seq2Seq Transformer Model",
        category="MODULE 10: LLMS & TRANSFORMERS",
        panel_title="Bilingual Machine Translation Model",
        bullets=[
            "Seq2Seq Architecture: Couples an Encoder (source comprehension) with a Decoder (target autoregressive generation).",
            "Pure Python Implementation: Standard tokenizer, embeddings, encoder self-attention, and decoder cross-attention blocks.",
            "Greedy Decoding Loop: Iterates forward passes, selecting the token with the maximum probability until a stop token is reached.",
            "Teacher Forcing: Feed actual ground truth targets into the decoder during training to stabilize parameters.",
            "Visual: Attention alignment weights mapping English to Spanish."
        ],
        img_path="plots/llm_5_attention_alignment.png",
        border_color=ACCENT_TEAL
    )

    # =============================================================
    # Slide 23: Module 11.1: Agentic AI Core Architecture & Concepts
    # =============================================================
    create_standard_slide(
        title="Agentic AI: Core Concepts & Architecture",
        category="MODULE 11: AGENTIC AI SYSTEMS",
        panel_title="Autonomous Reasoning Loops",
        bullets=[
            "Paradigm Shift: Moves from static text responses to dynamic reasoning agents acting as computer systems operators.",
            "Core Components:",
            "  - AI Agent: The controller managing history (memory) and execution loops.",
            "  - Thought-Action-Observation Loop: The reasoning cycle (e.g. ReAct).",
            "  - Atomic Tool: Low-level executable Python function or command.",
            "  - Composite Skill: Orchestrated multi-step workflow in code.",
            "  - Model Context Protocol (MCP): Open-standard communication layer.",
            "Function Calling: Structured JSON exchange between Agent and LLM."
        ],
        img_path="plots/agentic_concepts_diagram.png",
        border_color=ACCENT_TEAL
    )

    # =============================================================
    # Slide 24: Module 11.2: Runtime Introspection & Tool Schemas
    # =============================================================
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_slide_header(slide, "Runtime Introspection & Tool Schemas", "MODULE 11: AGENTIC AI SYSTEMS")
    
    # Left Panel: Explanation
    draw_panel(slide, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), border_color=ACCENT_INDIGO)
    box_left = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_left = box_left.text_frame
    tf_left.word_wrap = True
    
    p = tf_left.paragraphs[0]
    p.text = "Automated Schema Registration"
    p.font.name = "Arial"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(14)
    
    bullets_intro = [
        "Introspection Concept: Instead of manually writing error-prone JSON declarations, agent registries inspect code at runtime.",
        "Docstring Parsing: Scans descriptions using regular expressions to extract functional details and argument documentation.",
        "Signature Analysis: inspect.signature() inspects type annotations (e.g. str -> STRING, int -> INTEGER).",
        "Required vs Optional: Detects arguments lacking defaults to mark them as 'required' in the schema definition."
    ]
    for b in bullets_intro:
        p = tf_left.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_BODY
        p.space_after = Pt(10)
        
    # Right Panel: Code vs Schema
    draw_panel(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), border_color=ACCENT_ORANGE)
    box_right = slide.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.4))
    tf_right = box_right.text_frame
    tf_right.word_wrap = True
    
    p = tf_right.paragraphs[0]
    p.text = "Decorated Code to Schema Output"
    p.font.name = "Arial"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ORANGE
    p.space_after = Pt(14)
    
    code_text = (
        "# Python Function Definition:\n"
        "@tool\n"
        "def get_weather(city: str) -> str:\n"
        "    \"\"\"Fetches weather. Args: city: Target city\"\"\"\n"
        "    return f\"Weather in {city}: 18°C\"\n\n"
        "# Generated JSON Tool Schema:\n"
        "{\n"
        "  \"name\": \"get_weather\",\n"
        "  \"description\": \"Fetches weather...\",\n"
        "  \"parameters\": {\n"
        "    \"type\": \"OBJECT\",\n"
        "    \"properties\": {\n"
        "      \"city\": {\"type\": \"STRING\", \"description\": \"Target city\"}\n"
        "    },\n"
        "    \"required\": [\"city\"]\n"
        "  }\n"
        "}"
    )
    p_code = tf_right.add_paragraph()
    p_code.text = code_text
    p_code.font.name = "Courier New"
    p_code.font.size = Pt(11)
    p_code.font.bold = True
    p_code.font.color.rgb = TEXT_TITLE
    p_code.space_after = Pt(4)

    # =============================================================
    # Slide 25: Module 11.3: Composite Skills & Dynamic Code Execution
    # =============================================================
    create_standard_slide(
        title="Composite Skills & Dynamic Execution",
        category="MODULE 11: AGENTIC AI SYSTEMS",
        panel_title="Encapsulating Local Workflows",
        bullets=[
            "Composite Skill Concept: Packages complex multi-step tasks to run locally, avoiding multiple remote network API roundtrips.",
            "Skill Package Directory Layout:",
            "  - Directory name maps to skill (e.g. skills/research_city/).",
            "  - SKILL.md: YAML metadata defining schema parameters + instructions.",
            "  - script.py: Local Python execution orchestrating underlying tools.",
            "Dynamic Execution via exec():",
            "  - Load script.py contents at runtime.",
            "  - Inject local tools and kwargs into exec_globals/locals namespaces.",
            "  - exec(script_code, globals, locals) runs code dynamically.",
            "  - Capture and return the resulting value via result variable."
        ],
        img_path="plots/agentic_loop_sequence.png",
        border_color=ACCENT_PURPLE
    )

    # =============================================================
    # Slide 26: Module 11.4: Model Context Protocol (MCP)
    # =============================================================
    create_standard_slide(
        title="Model Context Protocol (MCP) Architecture",
        category="MODULE 11: AGENTIC AI SYSTEMS",
        panel_title="Standardizing Tool Integration Interfaces",
        bullets=[
            "MCP Architecture: Open standard that decouples tools providers (servers) from AI applications (clients).",
            "Three Core Actors in the Protocol:",
            "  1. MCP Client: Protocol consumer (e.g. IDE, AI agent host).",
            "  2. MCP Host: Runs the client session and orchestrates context.",
            "  3. MCP Server: Provides tools, resources, and prompt templates.",
            "Standardized Transport Mechanisms:",
            "  - Stdio Transport: JSON-RPC over stdin/stdout (same-machine).",
            "  - SSE Transport: Server-Sent Events stream + POST requests (network).",
            "Benefits: Single protocol connects any agent to files, DBs, and APIs."
        ],
        img_path="plots/mcp_architecture_diagram.png",
        border_color=ACCENT_ORANGE
    )

    # =============================================================
    # Slide 27: Appendix: Math Glossary
    # =============================================================
    create_standard_slide(
        title="Appendix: Centroid, Hyperplane & Residuals",
        category="APPENDIX: GLOSSARY",
        panel_title="Mathematical Fundamentals",
        bullets=[
            "Centroid: The geometric center of a cluster of coordinates. Computed as the mean vector: mu = 1/N * sum(x_i).",
            "Hyperplane: A linear boundary of dimension D-1 that separates a D-dimensional space: w^T * x + b = 0.",
            "Margin: The distance between separating hyperplane and the closest support vectors: 2 / ||w||.",
            "Residual: The difference between actual and predicted target values: e_i = y_i - y_hat_i.",
            "Visual: Side-by-side plots of Centroid coordinate mean, SVM margins/support-vectors, and OLS regression residuals."
        ],
        img_path="plots/glossary_math_concepts.png",
        border_color=ACCENT_INDIGO
    )

    # =============================================================
    # Slide 28: Appendix: Normalization & Regularization
    # =============================================================
    create_standard_slide(
        title="Appendix: Normalization & Regularization",
        category="APPENDIX: GLOSSARY",
        panel_title="Scaling and Penalty Geometry",
        bullets=[
            "Normalization: Rescaling features to a shared coordinate bounds.",
            "  - Min-Max Scaling: Maps features to a fixed [0, 1] range.",
            "  - Standardization: Transforms features to have mean=0 and std=1.",
            "Regularization: Constrains model parameter magnitudes to prevent overfitting.",
            "  - Lasso (L1): Adds absolute weight penalty. Creates diamond-shaped boundaries that drive weights to exactly 0 (sparsity).",
            "  - Ridge (L2): Adds squared weight penalty. Circular boundaries that shrink weights close to 0 but keep all features active."
        ],
        img_path="plots/glossary_normalization_regularization.png",
        border_color=ACCENT_TEAL
    )

    # =============================================================
    # Slide 29: Appendix: Optimization & Bias Glossary
    # =============================================================
    create_standard_slide(
        title="Appendix: Optimization & Bias",
        category="APPENDIX: GLOSSARY",
        panel_title="Gradients and Algorithmic Assumptions",
        bullets=[
            "Gradient: Vector of partial derivatives pointing in the direction of steepest rate of function increase: [df/dw_1, ..., df/dw_D].",
            "Gradient Descent: Updates weights along negative gradient: w = w - eta * grad.",
            "Adam Optimizer: Computes adaptive learning rates using first moment (momentum) and second moment (gradient variance). Scales updates per coordinate.",
            "Inductive Bias: Mathematical assumptions a model makes to generalize to unseen inputs (e.g., linear vs. locality bias).",
            "Lazy Learning: Defers generalization until query time (e.g. KNN). Features zero training time but slow prediction.",
            "Visual: Gradient 3D valley, linear/locality boundaries, KNN distance query."
        ],
        img_path="plots/glossary_gradient_bias_lazy.png",
        border_color=ACCENT_ORANGE
    )

    # =============================================================
    # Slide 30: Appendix: LLM Concepts Glossary
    # =============================================================
    create_standard_slide(
        title="Appendix: LLM Core Terminology",
        category="APPENDIX: GLOSSARY",
        panel_title="Transformer Building Blocks",
        bullets=[
            "Token: The basic subword unit mapping to vocabulary IDs.",
            "Self-Attention: Maps query (Q) alignments to keys (K) to weight values (V).",
            "Causal Mask: Restricts attention to past positions.",
            "Autoregressive Decoding: Appends output tokens back to input sequence.",
            "Temperature: Scaling factor applied to logits before softmax.",
            "  - Low Temp (T -> 0): Peakier probabilities, deterministic generation.",
            "  - High Temp (T -> infinity): Uniform probabilities, highly random/creative.",
            "Visual: Token IDs, Self-Attention nodes, Mask heatmap, and Temperature curves."
        ],
        img_path="plots/glossary_llm_concepts.png",
        border_color=ACCENT_PURPLE
    )

    # =============================================================
    # Slide 31: Course Summary & Teaching Strategy
    # =============================================================
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_slide_header(slide, "Teaching Strategy & Lecture Structure", "COURSE CONCLUSION")
    
    # Large center box
    draw_panel(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8), border_color=ACCENT_TEAL)
    box_sum = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.9), Inches(4.4))
    tf_sum = box_sum.text_frame
    tf_sum.word_wrap = True
    
    p = tf_sum.paragraphs[0]
    p.text = "Recommended Lecture Structure"
    p.font.name = "Arial"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(14)
    
    sum_pts = [
        "1. Start with Paradigms: Make students classify everyday examples manually.",
        "2. Regression vs. Classification: Differentiate numeric forecasting vs. discrete categorization.",
        "3. Mathematics to Code: Show the formulas (Lasso, Logistic) side-by-side with scikit-learn models.",
        "4. Emphasize Evaluation: A model scoring 99% accuracy on training data is almost always overfit.",
        "5. Hands-on labs: Have students tune K-Means centroids or Random Forest depths on standard datasets (Iris/Titanic).",
        "6. Introduce Generative AI: Transition from embeddings (Word2Vec) to sequential comprehension (Transformers).",
        "7. Explain Autonomous systems: Focus on reasoning loops (thought/action/observation) and standardized tooling interfaces (MCP)."
    ]
    for pt in sum_pts:
        p = tf_sum.add_paragraph()
        p.text = pt
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_BODY
        p.space_after = Pt(12)
        
    # Save Presentation
    output_filename = "machine_learning_fundamentals.pptx"
    prs.save(output_filename)
    print(f"Presentation successfully created and saved as: '{output_filename}'")

if __name__ == "__main__":
    create_deck()
